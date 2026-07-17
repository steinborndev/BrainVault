#!/usr/bin/env python3
"""Extract plain text from PowerPoint / Excel / OpenDocument files to stdout.

Used by the office preprocessing plugin for the formats pandoc cannot read
(SPEC.md §5). Dispatches on extension; prints a readable plaintext rendering.

Dependencies (pip, no sudo): python-pptx, openpyxl, odfpy.
Exits non-zero with a message on stderr if a required library is missing, so the
plugin surfaces "tool not installed" rather than a silent empty extraction.
"""
import sys
import os


def die(msg: str) -> "NoReturn":  # type: ignore[name-defined]
    print(msg, file=sys.stderr)
    sys.exit(1)


def extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        die("python-pptx is not installed (pip install python-pptx)")
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        out.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    out.append(" | ".join(cells))
        out.append("")
    return "\n".join(out)


def extract_xlsx(path: str) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        die("openpyxl is not installed (pip install openpyxl)")
    wb = load_workbook(path, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"## Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                out.append(" | ".join(cells))
        out.append("")
    return "\n".join(out)


def extract_odf(path: str) -> str:
    try:
        from odf.opendocument import load
        from odf import text as odf_text, teletype
    except ImportError:
        die("odfpy is not installed (pip install odfpy)")
    doc = load(path)
    out = []
    for para in doc.getElementsByType(odf_text.P):
        line = teletype.extractText(para).strip()
        if line:
            out.append(line)
    return "\n".join(out)


def main() -> None:
    if len(sys.argv) != 2:
        die("usage: extract-office.py <file>")
    path = sys.argv[1]
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    if ext in ("pptx", "ppt"):
        text = extract_pptx(path)
    elif ext in ("xlsx", "xls"):
        text = extract_xlsx(path)
    elif ext in ("ods", "odp", "odt"):
        text = extract_odf(path)
    else:
        die(f"unsupported office extension: {ext}")
    sys.stdout.write(text)


if __name__ == "__main__":
    main()
